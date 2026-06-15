from __future__ import annotations

from copy import deepcopy
from pathlib import Path
import json
import re

from pptx import Presentation
from pptx.slide import Slide

from pptx_gen.schemas import GeneratedElement, PageGenerationResult


class PPTBuilder:
    def __init__(self, template_path: str | Path) -> None:
        self.template_path = Path(template_path)

    def build(self, rules_path: str | Path, page_results_path: str | Path, output_path: str | Path) -> Path:
        rules = json.loads(Path(rules_path).read_text(encoding="utf-8"))
        page_results = self._load_page_results(page_results_path)
        page_result_map = {item.page_no: item for item in page_results}
        page_rule_map = {int(page["page_no"]): page for page in rules["pages"]}

        presentation = Presentation(self.template_path)
        for slide_index in reversed(range(len(presentation.slides))):
            page_no = slide_index + 1
            result = page_result_map.get(page_no)
            page_rule = page_rule_map.get(page_no)
            if result is None or page_rule is None or not result.should_generate:
                self._delete_slide(presentation, slide_index)
                continue
            split_results = self._expand_result_for_overflow(page_rule, result)
            slide = presentation.slides[slide_index]
            self._fill_slide(slide, page_rule, split_results[0])
            insert_after_index = slide_index
            for extra_result in split_results[1:]:
                duplicated_slide = self._duplicate_slide_after(presentation, insert_after_index)
                self._fill_slide(duplicated_slide, page_rule, extra_result)
                insert_after_index += 1

        target = Path(output_path)
        target.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(target)
        return target

    def _load_page_results(self, page_results_path: str | Path) -> list[PageGenerationResult]:
        root = Path(page_results_path)
        items: list[PageGenerationResult] = []
        for path in sorted(root.glob("page_*.json")):
            payload = json.loads(path.read_text(encoding="utf-8"))
            items.append(PageGenerationResult.model_validate(payload))
        return items

    def _fill_slide(self, slide: Slide, page_rule: dict, result: PageGenerationResult) -> None:
        rule_element_map = {item["id"]: item for item in page_rule.get("elements", [])}
        result_element_map = {item.id: item for item in result.elements}
        shape_map = {shape.shape_id: shape for shape in slide.shapes}
        for element_id, rule in rule_element_map.items():
            shape = shape_map.get(int(rule["shape_id"]))
            result_element = result_element_map.get(element_id)
            if rule["type"] in {"title", "text"}:
                if shape is not None and hasattr(shape, "text_frame"):
                    self._set_text(shape, (result_element.content if result_element is not None else "") or "")
            elif rule["type"] == "table":
                if shape is not None and getattr(shape, "has_table", False):
                    self._fill_table(shape.table, result_element)
            elif rule["type"] == "image":
                self._place_image(slide, shape, rule, result_element)

    def _expand_result_for_overflow(self, page_rule: dict, result: PageGenerationResult) -> list[PageGenerationResult]:
        if page_rule.get("page_purpose") != "text":
            return [result]
        if any(item.get("type") in {"table", "image"} for item in page_rule.get("elements", [])):
            return [result]

        result_element_map = {item.id: item for item in result.elements}
        body_rules = [
            item
            for item in page_rule.get("elements", [])
            if item.get("type") == "text"
            and item.get("id") in result_element_map
            and (result_element_map[item["id"]].content or "").strip()
        ]
        if len(body_rules) != 1:
            return [result]

        body_rule = body_rules[0]
        content_requirement = body_rule.get("content_requirement") or ""
        if "可分多页" not in content_requirement and "分多页" not in content_requirement:
            return [result]

        body_element = result_element_map[body_rule["id"]]
        chunks = self._split_text_to_fit(page_rule, body_rule, body_element.content or "")
        if len(chunks) <= 1:
            return [result]

        return [self._clone_result_with_text(result, body_element.id, chunk) for chunk in chunks]

    def _split_text_to_fit(self, page_rule: dict, rule: dict, text: str) -> list[str]:
        chars_per_line, line_slots = self._estimate_text_capacity(page_rule, rule)
        if chars_per_line <= 0 or line_slots <= 0:
            return [text]

        wrapped_blocks: list[list[str]] = []
        for raw_line in text.splitlines():
            stripped = raw_line.strip()
            if not stripped:
                continue
            wrapped = self._wrap_line(stripped, chars_per_line)
            if wrapped:
                wrapped_blocks.append(wrapped)

        wrapped_lines = [line for block in wrapped_blocks for line in block]

        if len(wrapped_lines) <= line_slots:
            return [text]

        chunks: list[str] = []
        current_lines: list[str] = []
        for block in wrapped_blocks:
            if len(block) > line_slots:
                if current_lines:
                    chunks.append("\n".join(current_lines).strip())
                    current_lines = []
                for start in range(0, len(block), line_slots):
                    part = block[start : start + line_slots]
                    chunks.append("\n".join(part).strip())
                continue

            if current_lines and len(current_lines) + len(block) > line_slots:
                chunks.append("\n".join(current_lines).strip())
                current_lines = list(block)
            else:
                current_lines.extend(block)
        if current_lines:
            chunks.append("\n".join(current_lines).strip())
        return [chunk for chunk in chunks if chunk]

    def _estimate_text_capacity(self, page_rule: dict, rule: dict) -> tuple[int, int]:
        bbox = rule.get("bbox") or {}
        style = rule.get("style") or {}
        margins = style.get("margins") or {}
        width_pt = max(int(bbox.get("w", 0)) - int(margins.get("left", 0)) - int(margins.get("right", 0)), 0) / 12700
        height_pt = max(int(bbox.get("h", 0)) - int(margins.get("top", 0)) - int(margins.get("bottom", 0)), 0) / 12700
        font_pt = max((style.get("font_size") or 1200) / 100, 10)
        font_hint_pt, line_spacing = self._extract_text_layout_hint(page_rule)
        font_pt = max(font_pt, font_hint_pt)
        chars_per_line = max(int(width_pt / (font_pt * 2.3)), 1)
        line_slots = max(int(height_pt / (font_pt * line_spacing * 1.15)), 1)
        return chars_per_line, line_slots

    def _extract_text_layout_hint(self, page_rule: dict) -> tuple[float, float]:
        font_pt = 0.0
        line_spacing = 1.4
        for element in page_rule.get("elements", []):
            if element.get("type") != "text":
                continue
            default_text = element.get("default_text") or ""
            font_match = re.search(r"(\d+(?:\.\d+)?)\s*磅", default_text)
            spacing_match = re.search(r"(\d+(?:\.\d+)?)\s*倍行间距", default_text)
            if font_match:
                font_pt = max(font_pt, float(font_match.group(1)))
            if spacing_match:
                line_spacing = max(line_spacing, float(spacing_match.group(1)))
        return font_pt, line_spacing

    def _wrap_line(self, text: str, width: int) -> list[str]:
        if len(text) <= width:
            return [text]

        wrapped: list[str] = []
        remaining = text
        while len(remaining) > width:
            split_at = self._find_split_position(remaining, width)
            wrapped.append(remaining[:split_at].rstrip())
            remaining = remaining[split_at:].lstrip()
        if remaining:
            wrapped.append(remaining)
        return wrapped

    def _find_split_position(self, text: str, width: int) -> int:
        candidate = min(width, len(text))
        lower_bound = max(width // 2, 1)
        while candidate > lower_bound:
            if re.match(r"[，。；：、,.;:）)]", text[candidate - 1]):
                return candidate
            candidate -= 1
        return min(width, len(text))

    def _clone_result_with_text(self, result: PageGenerationResult, element_id: str, content: str) -> PageGenerationResult:
        cloned_elements: list[GeneratedElement] = []
        for element in result.elements:
            cloned = element.model_copy(deep=True)
            if cloned.id == element_id:
                cloned.content = content
            cloned_elements.append(cloned)
        return PageGenerationResult(
            page_no=result.page_no,
            should_generate=result.should_generate,
            skip_reason=result.skip_reason,
            elements=cloned_elements,
        )

    def _duplicate_slide_after(self, presentation: Presentation, slide_index: int) -> Slide:
        source_slide = presentation.slides[slide_index]
        duplicated_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        for shape in list(duplicated_slide.shapes):
            self._remove_shape(shape)
        for shape in source_slide.shapes:
            duplicated_slide.shapes._spTree.insert_element_before(deepcopy(shape._element), "p:extLst")
        slide_id_list = presentation.slides._sldIdLst
        new_slide_id = slide_id_list[-1]
        del slide_id_list[-1]
        slide_id_list.insert(slide_index + 1, new_slide_id)
        return presentation.slides[slide_index + 1]

    def _set_text(self, shape, content: str) -> None:
        text_frame = shape.text_frame
        paragraph_alignment = None
        font_name = None
        font_size = None
        bold = None
        italic = None
        if text_frame.paragraphs:
            paragraph_alignment = text_frame.paragraphs[0].alignment
            if text_frame.paragraphs[0].runs:
                font = text_frame.paragraphs[0].runs[0].font
                font_name = font.name
                font_size = font.size
                bold = font.bold
                italic = font.italic
        text_frame.clear()
        paragraph = text_frame.paragraphs[0]
        paragraph.text = content
        paragraph.alignment = paragraph_alignment
        if paragraph.runs:
            font = paragraph.runs[0].font
            font.name = font_name
            font.size = font_size
            font.bold = bold
            font.italic = italic

    def _fill_table(self, table, element: GeneratedElement | None) -> None:
        all_rows = []
        if element is not None and element.headers:
            all_rows.append(element.headers)
        if element is not None and element.rows:
            all_rows.extend(element.rows)
        for row_index in range(len(table.rows)):
            for col_index in range(len(table.columns)):
                value = ""
                if row_index < len(all_rows) and col_index < len(all_rows[row_index]):
                    value = str(all_rows[row_index][col_index])
                table.cell(row_index, col_index).text = value

    def _place_image(self, slide: Slide, shape, rule: dict, element: GeneratedElement | None) -> None:
        if element is None or not element.rendered_path:
            if shape is not None:
                self._remove_shape(shape)
            return
        image_path = Path(element.rendered_path)
        if not image_path.exists():
            if shape is not None:
                self._remove_shape(shape)
            return
        if shape is not None:
            self._remove_shape(shape)
        bbox = rule["bbox"]
        slide.shapes.add_picture(
            str(image_path),
            left=int(bbox["x"]),
            top=int(bbox["y"]),
            width=int(bbox["w"]),
            height=int(bbox["h"]),
        )

    def _remove_shape(self, shape) -> None:
        element = shape._element
        parent = element.getparent()
        if parent is not None:
            parent.remove(element)

    def _delete_slide(self, presentation: Presentation, slide_index: int) -> None:
        slide_id_list = presentation.slides._sldIdLst
        slide_id = slide_id_list[slide_index]
        relationship_id = slide_id.rId
        presentation.part.drop_rel(relationship_id)
        del slide_id_list[slide_index]
